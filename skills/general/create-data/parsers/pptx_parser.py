"""PPTX text and slide structure extraction using python-pptx.

Extracts text content, titles, body text, and table data per slide,
output as structured markdown.

Usage:
    python pptx_parser.py <path_to_pptx>

Dependencies:
    pip install python-pptx
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches  # noqa: F401
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx")
    sys.exit(1)


def _extract_table_as_markdown(table) -> str:
    """Convert a PPTX table shape to a markdown table."""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append(cells)

    if not rows:
        return ""

    lines = []
    lines.append("| " + " | ".join(rows[0]) + " |")
    lines.append("| " + " | ".join("---" for _ in rows[0]) + " |")
    for row in rows[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def parse_pptx(pptx_path: str | Path) -> str:
    """Parse a PPTX file and return slide-by-slide structured markdown.

    Args:
        pptx_path: Path to the PPTX file.

    Returns:
        Structured markdown string with slide titles, body text, and tables.
    """
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX not found: {pptx_path}")

    prs = Presentation(str(pptx_path))
    sections = []

    for i, slide in enumerate(prs.slides, 1):
        slide_parts = [f"## Slide {i}"]

        # Extract title
        if slide.shapes.title and slide.shapes.title.text.strip():
            slide_parts.append(f"**Title:** {slide.shapes.title.text.strip()}")

        # Process shapes
        body_texts = []
        tables = []

        for shape in slide.shapes:
            # Skip title shape (already handled)
            if shape == slide.shapes.title:
                continue

            if shape.has_table:
                tables.append(_extract_table_as_markdown(shape.table))
            elif shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    body_texts.append(text)

        if body_texts:
            slide_parts.append("\n".join(body_texts))

        for table in tables:
            slide_parts.append(table)

        sections.append("\n\n".join(slide_parts))

    return f"# {pptx_path.name}\n\n" + "\n\n---\n\n".join(sections)


def main():
    parser = argparse.ArgumentParser(description="Extract structured content from PPTX files")
    parser.add_argument("pptx_path", help="Path to the PPTX file")
    args = parser.parse_args()

    result = parse_pptx(args.pptx_path)
    print(result)


if __name__ == "__main__":
    main()
